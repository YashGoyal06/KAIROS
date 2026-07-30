import pytest
import pytest_asyncio
import uuid
import io
from httpx import AsyncClient, ASGITransport
from sqlalchemy.ext.asyncio import create_async_engine, AsyncSession
from sqlalchemy.orm import sessionmaker
from pptx import Presentation

from backend.main import app
from backend.app.db.connection import get_db
from backend.app.db.models import Base, Profile, Session, Task, Team

# SQLite in-memory test database for fast, isolated test execution
TEST_DATABASE_URL = "sqlite+aiosqlite:///:memory:"

test_engine = create_async_engine(
    TEST_DATABASE_URL,
    echo=False,
    future=True
)

TestAsyncSessionLocal = sessionmaker(
    bind=test_engine,
    class_=AsyncSession,
    expire_on_commit=False,
    autocommit=False,
    autoflush=False
)

async def override_get_db():
    async with TestAsyncSessionLocal() as session:
        try:
            yield session
            await session.commit()
        except Exception:
            await session.rollback()
            raise
        finally:
            await session.close()

app.dependency_overrides[get_db] = override_get_db

@pytest_asyncio.fixture(autouse=True)
async def setup_database():
    async with test_engine.begin() as conn:
        await conn.run_sync(Base.metadata.create_all)
    yield
    async with test_engine.begin() as conn:
        await conn.run_sync(Base.metadata.drop_all)

@pytest_asyncio.fixture
async def async_client():
    transport = ASGITransport(app=app)
    async with AsyncClient(transport=transport, base_url="http://test") as client:
        yield client

@pytest.mark.asyncio
async def test_root_endpoint(async_client: AsyncClient):
    response = await async_client.get("/")
    assert response.status_code == 200
    assert response.json() == {"status": "online", "message": "KAIROS Core API"}

@pytest.mark.asyncio
async def test_templates_endpoint(async_client: AsyncClient):
    response = await async_client.get("/api/v1/sessions/00000000-0000-0000-0000-000000000000/pitch/templates")
    assert response.status_code == 200
    data = response.json()
    assert "templates" in data
    assert len(data["templates"]) == 5
    template_ids = [t["id"] for t in data["templates"]]
    assert "template-1" in template_ids
    assert "template-5" in template_ids

@pytest.mark.asyncio
async def test_end_to_end_pitch_pptx_export(async_client: AsyncClient):
    # 1. Create dummy profile
    profile_id = str(uuid.uuid4())
    prof_payload = {
        "id": profile_id,
        "full_name": "Test Architect",
        "primary_role": "Full Stack Developer",
        "experience_level": "Advanced",
        "tech_stack": ["Python", "FastAPI", "React", "PowerPoint Engine"]
    }
    prof_resp = await async_client.post("/api/v1/profiles", json=prof_payload)
    assert prof_resp.status_code in [200, 201]

    # 2. Create dummy session
    sess_payload = {
        "name": "KAIROS Presentation Test Deck",
        "creator_id": profile_id
    }
    sess_resp = await async_client.post("/api/v1/sessions", json=sess_payload)
    assert sess_resp.status_code in [200, 201]
    session_id = sess_resp.json()["id"]

    # 3. Export PPTX file for all 5 templates
    for tid in ["template-1", "template-2", "template-3", "template-4", "template-5"]:
        export_resp = await async_client.post(
            f"/api/v1/sessions/{session_id}/pitch/export-pptx",
            data={"template_id": tid}
        )
        assert export_resp.status_code == 200
        assert export_resp.headers["content-type"] == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        
        # Verify valid PPTX binary payload using pptx Presentation parser
        prs = Presentation(io.BytesIO(export_resp.content))
        assert len(prs.slides) > 0

@pytest.mark.asyncio
async def test_preview_slides_endpoint(async_client: AsyncClient):
    profile_id = str(uuid.uuid4())
    
    await async_client.post("/api/v1/profiles", json={
        "id": profile_id,
        "full_name": "Preview Tester",
        "primary_role": "AI Engineer",
        "experience_level": "Intermediate",
        "tech_stack": ["PyTorch", "FastAPI"]
    })
    sess_resp = await async_client.post("/api/v1/sessions", json={
        "name": "Preview Demo Project",
        "creator_id": profile_id
    })
    assert sess_resp.status_code in [200, 201]
    session_id = sess_resp.json()["id"]

    preview_resp = await async_client.post(
        f"/api/v1/sessions/{session_id}/pitch/preview-slides",
        data={"template_id": "template-1"}
    )
    assert preview_resp.status_code == 200
    data = preview_resp.json()
    assert "slides" in data
    assert len(data["slides"]) == 10
    assert isinstance(data["slides"][0], str)

@pytest.mark.asyncio
async def test_custom_template_analysis(async_client: AsyncClient):
    session_id = str(uuid.uuid4())
    
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Custom User Presentation Tag {PROJECT_NAME}"
    
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    
    response = await async_client.post(
        f"/api/v1/sessions/{session_id}/pitch/analyze-custom-template",
        files={"file": ("custom_test.pptx", buf.getvalue(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
    )
    assert response.status_code == 200
    res_data = response.json()
    assert res_data["status"] == "success"
    assert res_data["filename"] == "custom_test.pptx"
    assert res_data["analysis"]["slide_count"] == 1
