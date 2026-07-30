import os
import io
import re
import logging
from typing import Dict, Any, List, Optional
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from reportlab.lib.pagesizes import letter, landscape
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

logger = logging.getLogger("kairos.ppt_engine")

TEMPLATES_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "static", "ppt_templates")

PREDEFINED_TEMPLATES = {
    "template-gamma": {
        "id": "template-gamma",
        "name": "Gamma AI Ultra Studio",
        "file": "gamma",
        "description": "Gamma.app style 16:9 widescreen presentation with obsidian card containers, pill badges, and metric callouts.",
        "slides_count": 10
    },
    "template-1": {
        "id": "template-1",
        "name": "Cyber Neon Executive",
        "file": "Template-1.pptx",
        "description": "High-impact dark theme with neon purple accents. Ideal for AI & Tech Hackathons.",
        "slides_count": 10
    },
    "template-2": {
        "id": "template-2",
        "name": "Minimalist Modern Tech",
        "file": "Template-2.pptx",
        "description": "Clean, structured slide deck with high legibility and sleek grid layouts.",
        "slides_count": 11
    },
    "template-3": {
        "id": "template-3",
        "name": "Vibrant Launchpad",
        "file": "Template-3.pptx",
        "description": "Dynamic layout designed for pitch competitions, featuring prominent metric boxes.",
        "slides_count": 11
    },
    "template-4": {
        "id": "template-4",
        "name": "Enterprise Architecture",
        "file": "Template-4.pptx",
        "description": "Comprehensive design focusing on workflow, technical specs, and milestone flowcharts.",
        "slides_count": 13
    },
    "template-5": {
        "id": "template-5",
        "name": "Futuristic AI Studio",
        "file": "Template-5.pptx",
        "description": "Sleek obsidian gradient deck highlighting AI solution capabilities and team power.",
        "slides_count": 10
    }
}


class PPTEngine:

    @staticmethod
    def get_template_path(template_id: str) -> str:
        t_info = PREDEFINED_TEMPLATES.get(template_id, PREDEFINED_TEMPLATES["template-gamma"])
        return os.path.join(TEMPLATES_DIR, t_info["file"])

    @staticmethod
    def analyze_presentation(prs: Presentation) -> Dict[str, Any]:
        slide_details = []
        total_shapes = 0
        total_text_frames = 0

        for idx, slide in enumerate(prs.slides):
            text_boxes = []
            for shape in slide.shapes:
                total_shapes += 1
                if shape.has_text_frame:
                    total_text_frames += 1
                    text_content = shape.text_frame.text.strip()
                    if text_content:
                        text_boxes.append({
                            "shape_name": shape.name,
                            "char_count": len(text_content),
                            "sample_text": text_content[:60]
                        })
            
            slide_details.append({
                "slide_number": idx + 1,
                "shape_count": len(slide.shapes),
                "text_box_count": len(text_boxes),
                "sample_content": text_boxes[:3]
            })

        return {
            "slide_count": len(prs.slides),
            "total_shapes": total_shapes,
            "total_text_frames": total_text_frames,
            "slide_width_inches": round(prs.slide_width.inches, 2),
            "slide_height_inches": round(prs.slide_height.inches, 2),
            "slide_details": slide_details
        }

    @staticmethod
    def fit_text_to_frame(text_frame, text: str, max_font_size: int = 16, min_font_size: int = 9, is_title: bool = False):
        text_frame.word_wrap = True
        cleaned = text.strip()
        if not cleaned:
            return

        char_len = len(cleaned)
        
        if is_title:
            if char_len > 80:
                font_size = max(14, max_font_size - 6)
                cleaned = cleaned[:100] + "..." if char_len > 100 else cleaned
            elif char_len > 40:
                font_size = max(16, max_font_size - 4)
            else:
                font_size = max_font_size
        else:
            if char_len > 400:
                font_size = min_font_size
                cleaned = cleaned[:450] + "..."
            elif char_len > 250:
                font_size = max(min_font_size, max_font_size - 5)
            elif char_len > 140:
                font_size = max(min_font_size + 1, max_font_size - 3)
            else:
                font_size = max_font_size

        orig_color = None
        orig_name = None
        orig_bold = None
        orig_italic = None
        orig_alignment = None

        if len(text_frame.paragraphs) > 0:
            orig_alignment = text_frame.paragraphs[0].alignment
            if len(text_frame.paragraphs[0].runs) > 0:
                r = text_frame.paragraphs[0].runs[0]
                if r.font and r.font.color and r.font.color.type == 1:
                    orig_color = r.font.color.rgb
                if r.font and r.font.name:
                    orig_name = r.font.name
                if r.font and r.font.bold is not None:
                    orig_bold = r.font.bold
                if r.font and r.font.italic is not None:
                    orig_italic = r.font.italic

        text_frame.clear()
        
        lines = [line.strip() for line in cleaned.split("\n") if line.strip()]
        if not lines:
            lines = [cleaned]

        for i, line in enumerate(lines):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            if orig_alignment is not None:
                p.alignment = orig_alignment

            if line.startswith(("- ", "* ", "• ")):
                p.text = line[2:].strip()
                p.level = 0
            else:
                p.text = line

            p.font.size = Pt(font_size)
            if orig_color:
                p.font.color.rgb = orig_color
            if orig_name:
                p.font.name = orig_name
            if orig_bold is not None:
                p.font.bold = orig_bold
            if orig_italic is not None:
                p.font.italic = orig_italic

    @staticmethod
    def build_gamma_presentation(
        session_name: str,
        problem_statement: str,
        user_idea: str,
        pitch_sections: Dict[str, str],
        milestones: List[Dict[str, Any]] = None,
        tasks: List[Dict[str, Any]] = None,
        team_data: Dict[str, Any] = None
    ) -> bytes:
        """
        Gamma.app Style Presentation Generator:
        Generates 16:9 widescreen dark obsidian slides with card containers,
        pill badges, metric callout blocks, and crisp typography.
        """
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]  # Blank layout

        # Colors Palette (Gamma Obsidian Theme)
        BG_COLOR = RGBColor(9, 10, 15)          # Dark Obsidian #090a0f
        CARD_BG = RGBColor(18, 20, 30)          # Card Fill #12141e
        CARD_BORDER = RGBColor(42, 36, 56)      # Subtle Border #2a2438
        ACCENT_PURPLE = RGBColor(168, 85, 247)  # #a855f7
        ACCENT_CYAN = RGBColor(56, 189, 248)    # #38bdf8
        ACCENT_EMERALD = RGBColor(52, 211, 153) # #34d399
        TEXT_WHITE = RGBColor(255, 255, 255)
        TEXT_MUTED = RGBColor(148, 163, 184)    # #94a3b8

        team_info = team_data or {}
        user_name = team_info.get('name', 'Innovator')
        user_role = team_info.get('role', 'Fullstack Engineer')
        skills = team_info.get('skills', ['Python', 'React', 'FastAPI'])
        skills_str = ', '.join(skills) if isinstance(skills, list) else str(skills)

        def add_background(slide):
            bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = BG_COLOR
            bg_shape.line.fill.background()

        def add_badge(slide, text: str, accent_color=ACCENT_PURPLE, left=Inches(0.8), top=Inches(0.6)):
            badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.2), Inches(0.4))
            badge.fill.solid()
            badge.fill.fore_color.rgb = accent_color
            badge.line.fill.background()
            tf = badge.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = text.upper()
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = TEXT_WHITE
            p.font.name = 'Outfit'

        def add_card(slide, left, top, width, height, border_color=CARD_BORDER, bg_color=CARD_BG):
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
            card.fill.solid()
            card.fill.fore_color.rgb = bg_color
            card.line.color.rgb = border_color
            card.line.width = Pt(1.5)
            return card

        # Slide 1: Hero Cover Slide
        s1 = prs.slides.add_slide(blank_layout)
        add_background(s1)
        add_badge(s1, "KAIROS EXECUTION ENGINE", ACCENT_PURPLE, Inches(0.8), Inches(0.8))

        # Title
        t_box = s1.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.5))
        tf = t_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = session_name or "AI Hackathon Co-Founder"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.font.name = 'Outfit'

        p2 = tf.add_paragraph()
        p2.text = user_idea or "Autonomous AI-Driven Project Planning & Pitch Studio Platform"
        p2.font.size = Pt(18)
        p2.font.color.rgb = TEXT_MUTED
        p2.font.name = 'Inter'
        p2.space_before = Pt(8)

        # 3 Hero Feature Cards at Bottom
        card_w = Inches(3.64)
        card_gap = Inches(0.38)
        c1_bullets = [
            ("01. AI Coach Engine", "Real-time AI architecture & scope advisor."),
            ("02. Task Board Sync", "Automated priority task breakdown."),
            ("03. Gamma Studio Export", "Anti-overflow PPTX presentation suite.")
        ]
        for idx, (ctitle, cdesc) in enumerate(c1_bullets):
            c_left = Inches(0.8) + idx * (card_w + card_gap)
            c = add_card(s1, c_left, Inches(3.8), card_w, Inches(2.6), ACCENT_PURPLE if idx == 0 else CARD_BORDER)
            tf_c = c.text_frame
            tf_c.word_wrap = True
            p_ct = tf_c.paragraphs[0]
            p_ct.text = ctitle
            p_ct.font.size = Pt(16)
            p_ct.font.bold = True
            p_ct.font.color.rgb = ACCENT_PURPLE if idx == 0 else TEXT_WHITE
            p_ct.font.name = 'Outfit'

            p_cd = tf_c.add_paragraph()
            p_cd.text = cdesc
            p_cd.font.size = Pt(13)
            p_cd.font.color.rgb = TEXT_MUTED
            p_cd.font.name = 'Inter'
            p_cd.space_before = Pt(10)

        # Presenter Tag at Footer
        f_box = s1.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.5))
        p_f = f_box.text_frame.paragraphs[0]
        p_f.text = f"Presented By: {user_name} ({user_role})  •  Hackathon Showcase 2026"
        p_f.font.size = Pt(11)
        p_f.font.color.rgb = ACCENT_CYAN
        p_f.font.name = 'Inter'

        # Slide 2: Problem Statement (3 Cards Layout)
        s2 = prs.slides.add_slide(blank_layout)
        add_background(s2)
        add_badge(s2, "THE CHALLENGE", ACCENT_CYAN, Inches(0.8), Inches(0.6))

        t2_box = s2.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.7), Inches(1.0))
        p2_t = t2_box.text_frame.paragraphs[0]
        p2_t.text = "Problem Statement & Market Friction"
        p2_t.font.size = Pt(28)
        p2_t.font.bold = True
        p2_t.font.color.rgb = TEXT_WHITE
        p2_t.font.name = 'Outfit'

        p2_sub = t2_box.text_frame.add_paragraph()
        p2_sub.text = problem_statement or "Developers lose pitch momentum due to manual deck formatting and scope creep."
        p2_sub.font.size = Pt(14)
        p2_sub.font.color.rgb = TEXT_MUTED

        prob_cards = [
            ("Loss of Momentum", "Hackathon teams waste up to 40% of time formatting slides instead of building features."),
            ("Scope Creep & Blockers", "Unstructured tasks lead to missed deadlines during crunch hours."),
            ("Static Text Decks", "Traditional PPT generators output plain text templates that fail to impress judges.")
        ]
        for idx, (title_p, desc_p) in enumerate(prob_cards):
            c_left = Inches(0.8) + idx * (card_w + card_gap)
            c = add_card(s2, c_left, Inches(2.6), card_w, Inches(4.0))
            tf_c = c.text_frame
            tf_c.word_wrap = True
            
            p_h = tf_c.paragraphs[0]
            p_h.text = f"0{idx+1}."
            p_h.font.size = Pt(20)
            p_h.font.bold = True
            p_h.font.color.rgb = ACCENT_CYAN

            p_t = tf_c.add_paragraph()
            p_t.text = title_p
            p_t.font.size = Pt(18)
            p_t.font.bold = True
            p_t.font.color.rgb = TEXT_WHITE
            p_t.space_before = Pt(8)

            p_d = tf_c.add_paragraph()
            p_d.text = desc_p
            p_d.font.size = Pt(13)
            p_d.font.color.rgb = TEXT_MUTED
            p_d.space_before = Pt(12)

        # Slide 3: Core Solution (Hero Split Layout)
        s3 = prs.slides.add_slide(blank_layout)
        add_background(s3)
        add_badge(s3, "CORE SOLUTION", ACCENT_EMERALD, Inches(0.8), Inches(0.6))

        t3_box = s3.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.8))
        p3_t = t3_box.text_frame.paragraphs[0]
        p3_t.text = "The KAIROS Execution Engine"
        p3_t.font.size = Pt(28)
        p3_t.font.bold = True
        p3_t.font.color.rgb = TEXT_WHITE

        # Left Metric Callout Box
        c_left_box = add_card(s3, Inches(0.8), Inches(2.2), Inches(4.5), Inches(4.5), ACCENT_EMERALD)
        tf_l = c_left_box.text_frame
        tf_l.word_wrap = True

        p_m1 = tf_l.paragraphs[0]
        p_m1.text = "10x"
        p_m1.font.size = Pt(64)
        p_m1.font.bold = True
        p_m1.font.color.rgb = ACCENT_EMERALD

        p_m2 = tf_l.add_paragraph()
        p_m2.text = "Faster Presentation Generation"
        p_m2.font.size = Pt(20)
        p_m2.font.bold = True
        p_m2.font.color.rgb = TEXT_WHITE

        p_m3 = tf_l.add_paragraph()
        p_m3.text = pitch_sections.get("demo", "Real-time AI workflow engine for execution teams.")[:140]
        p_m3.font.size = Pt(13)
        p_m3.font.color.rgb = TEXT_MUTED
        p_m3.space_before = Pt(12)

        # Right Solution Features List Cards
        right_w = Inches(6.8)
        c_right = add_card(s3, Inches(5.7), Inches(2.2), right_w, Inches(4.5))
        tf_r = c_right.text_frame
        tf_r.word_wrap = True

        features = [
            ("Interactive AI Coach Room", "Generates execution roadmaps and resolves technical blockers on demand."),
            ("Task Board & Supabase Sync", "Real-time task priority tracking and milestone progress updates."),
            ("Gamma Ultra Slide Engine", "Guarantees zero text overflow and pristine vector card layouts.")
        ]
        for i, (ft, fd) in enumerate(features):
            p_ft = tf_r.paragraphs[0] if i == 0 else tf_r.add_paragraph()
            p_ft.text = f"• {ft}"
            p_ft.font.size = Pt(16)
            p_ft.font.bold = True
            p_ft.font.color.rgb = ACCENT_PURPLE
            if i > 0:
                p_ft.space_before = Pt(16)

            p_fd = tf_r.add_paragraph()
            p_fd.text = fd
            p_fd.font.size = Pt(13)
            p_fd.font.color.rgb = TEXT_MUTED
            p_fd.space_before = Pt(4)

        # Slide 4: System Architecture (4 Grid Cards)
        s4 = prs.slides.add_slide(blank_layout)
        add_background(s4)
        add_badge(s4, "SYSTEM ARCHITECTURE", ACCENT_PURPLE, Inches(0.8), Inches(0.6))

        t4_box = s4.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.8))
        p4_t = t4_box.text_frame.paragraphs[0]
        p4_t.text = "High-Performance Tech Stack"
        p4_t.font.size = Pt(28)
        p4_t.font.bold = True
        p4_t.font.color.rgb = TEXT_WHITE

        arch_cards = [
            ("FastAPI Backend", "Async Python server with SQLAlchemy ORM and SSE streaming."),
            ("Supabase DB & Auth", "PostgreSQL database with real-time subscriptions and secure JWT auth."),
            ("LLM Orchestrator", "Multi-model fallback broker routing Groq, Nvidia NIM, and Hugging Face."),
            ("React 19 & Vite UI", "Modern glassmorphism interface with Lenis smooth scroll and Tailwind CSS.")
        ]
        grid_w = Inches(5.6)
        grid_h = Inches(2.0)
        positions = [
            (Inches(0.8), Inches(2.2)),
            (Inches(6.9), Inches(2.2)),
            (Inches(0.8), Inches(4.6)),
            (Inches(6.9), Inches(4.6))
        ]
        for idx, (atitle, adesc) in enumerate(arch_cards):
            gx, gy = positions[idx]
            gc = add_card(s4, gx, gy, grid_w, grid_h)
            gtf = gc.text_frame
            gtf.word_wrap = True
            
            gp1 = gtf.paragraphs[0]
            gp1.text = atitle
            gp1.font.size = Pt(17)
            gp1.font.bold = True
            gp1.font.color.rgb = ACCENT_CYAN if idx % 2 == 0 else ACCENT_PURPLE

            gp2 = gtf.add_paragraph()
            gp2.text = adesc
            gp2.font.size = Pt(13)
            gp2.font.color.rgb = TEXT_MUTED
            gp2.space_before = Pt(8)

        # Slide 5: Roadmap & Execution (Step Pipeline)
        s5 = prs.slides.add_slide(blank_layout)
        add_background(s5)
        add_badge(s5, "ROADMAP & MILESTONES", ACCENT_CYAN, Inches(0.8), Inches(0.6))

        t5_box = s5.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.8))
        p5_t = t5_box.text_frame.paragraphs[0]
        p5_t.text = "Sprint Execution Milestones"
        p5_t.font.size = Pt(28)
        p5_t.font.bold = True
        p5_t.font.color.rgb = TEXT_WHITE

        ms_list = milestones or [
            {"title": "Phase 1: Architecture & DB Setup", "status": "Completed"},
            {"title": "Phase 2: AI Coach & Task Sync", "status": "Completed"},
            {"title": "Phase 3: Gamma Presentation Studio", "status": "In Progress"}
        ]
        ms_w = Inches(3.64)
        for idx, m in enumerate(ms_list[:3]):
            m_left = Inches(0.8) + idx * (ms_w + card_gap)
            mc = add_card(s5, m_left, Inches(2.5), ms_w, Inches(4.0))
            mtf = mc.text_frame
            mtf.word_wrap = True

            mp1 = mtf.paragraphs[0]
            mp1.text = f"STAGE 0{idx+1}"
            mp1.font.size = Pt(12)
            mp1.font.bold = True
            mp1.font.color.rgb = ACCENT_EMERALD

            mp2 = mtf.add_paragraph()
            mp2.text = m.get('title', 'Milestone')
            mp2.font.size = Pt(18)
            mp2.font.bold = True
            mp2.font.color.rgb = TEXT_WHITE
            mp2.space_before = Pt(8)

            mp3 = mtf.add_paragraph()
            mp3.text = f"Status: {m.get('status', 'In Progress')}"
            mp3.font.size = Pt(13)
            mp3.font.color.rgb = ACCENT_CYAN
            mp3.space_before = Pt(16)

        # Slide 6: Sprint Tasks & Priorities
        s6 = prs.slides.add_slide(blank_layout)
        add_background(s6)
        add_badge(s6, "SPRINT METRICS", ACCENT_EMERALD, Inches(0.8), Inches(0.6))

        t6_box = s6.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.8))
        p6_t = t6_box.text_frame.paragraphs[0]
        p6_t.text = "Real-Time Execution Task Board"
        p6_t.font.size = Pt(28)
        p6_t.font.bold = True
        p6_t.font.color.rgb = TEXT_WHITE

        t_list = tasks[:6] if tasks else [
            {"name": "FastAPI Endpoint Integration", "priority": "High"},
            {"name": "Gamma AI Slide Engine", "priority": "High"},
            {"name": "Supabase Connection Pool", "priority": "Medium"}
        ]
        t_card = add_card(s6, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5))
        ttf = t_card.text_frame
        ttf.word_wrap = True

        for i, titem in enumerate(t_list):
            tp = ttf.paragraphs[0] if i == 0 else ttf.add_paragraph()
            tp.text = f"• {titem.get('name', 'Task')}  [{titem.get('priority', 'High')} Priority]"
            tp.font.size = Pt(16)
            tp.font.bold = True
            tp.font.color.rgb = ACCENT_PURPLE if i % 2 == 0 else ACCENT_CYAN
            if i > 0:
                tp.space_before = Pt(14)

        # Slide 7: Team & Capabilities
        s7 = prs.slides.add_slide(blank_layout)
        add_background(s7)
        add_badge(s7, "OUR TEAM", ACCENT_PURPLE, Inches(0.8), Inches(0.6))

        t7_box = s7.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.8))
        p7_t = t7_box.text_frame.paragraphs[0]
        p7_t.text = "Technical Mastery & Innovators"
        p7_t.font.size = Pt(28)
        p7_t.font.bold = True
        p7_t.font.color.rgb = TEXT_WHITE

        # Left Member Card
        m_card = add_card(s7, Inches(0.8), Inches(2.2), Inches(5.6), Inches(4.5))
        mtf = m_card.text_frame
        mtf.word_wrap = True

        mp1 = mtf.paragraphs[0]
        mp1.text = user_name
        mp1.font.size = Pt(24)
        mp1.font.bold = True
        mp1.font.color.rgb = TEXT_WHITE

        mp2 = mtf.add_paragraph()
        mp2.text = user_role
        mp2.font.size = Pt(16)
        mp2.font.color.rgb = ACCENT_PURPLE
        mp2.space_before = Pt(4)

        mp3 = mtf.add_paragraph()
        mp3.text = f"Core Technical Stack:\n{skills_str}"
        mp3.font.size = Pt(14)
        mp3.font.color.rgb = TEXT_MUTED
        mp3.space_before = Pt(16)

        # Right Capability Card
        cap_card = add_card(s7, Inches(6.9), Inches(2.2), Inches(5.6), Inches(4.5))
        ctf = cap_card.text_frame
        ctf.word_wrap = True

        cp1 = ctf.paragraphs[0]
        cp1.text = "Engineering Capabilities"
        cp1.font.size = Pt(20)
        cp1.font.bold = True
        cp1.font.color.rgb = ACCENT_CYAN

        capabilities = [
            "Fullstack Async Architecture (FastAPI + React 19)",
            "AI Agent Orchestration & Multi-LLM Cascading",
            "Zero-Overflow Vector Slide Generation"
        ]
        for cap in capabilities:
            cp = ctf.add_paragraph()
            cp.text = f"✔ {cap}"
            cp.font.size = Pt(14)
            cp.font.color.rgb = TEXT_WHITE
            cp.space_before = Pt(12)

        # Slide 8: Pitch Showcase Impact
        s8 = prs.slides.add_slide(blank_layout)
        add_background(s8)
        add_badge(s8, "VALUE PROPOSITION", ACCENT_EMERALD, Inches(0.8), Inches(0.6))

        t8_box = s8.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.8))
        p8_t = t8_box.text_frame.paragraphs[0]
        p8_t.text = "Pitch Showcase & Key Insight"
        p8_t.font.size = Pt(28)
        p8_t.font.bold = True
        p8_t.font.color.rgb = TEXT_WHITE

        q_card = add_card(s8, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5), ACCENT_EMERALD)
        qtf = q_card.text_frame
        qtf.word_wrap = True

        qp1 = qtf.paragraphs[0]
        qp1.text = f"“ {pitch_sections.get('showcase', 'KAIROS provides an end-to-end execution co-founder, turning hackathon ideas into production software rapidly.')} ”"
        qp1.font.size = Pt(22)
        qp1.font.bold = True
        qp1.font.color.rgb = TEXT_WHITE

        # Slide 9: Risk Mitigation
        s9 = prs.slides.add_slide(blank_layout)
        add_background(s9)
        add_badge(s9, "RISK STRATEGY", ACCENT_CYAN, Inches(0.8), Inches(0.6))

        t9_box = s9.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.8))
        p9_t = t9_box.text_frame.paragraphs[0]
        p9_t.text = "Risk Management & Resiliency"
        p9_t.font.size = Pt(28)
        p9_t.font.bold = True
        p9_t.font.color.rgb = TEXT_WHITE

        r_left = add_card(s9, Inches(0.8), Inches(2.2), Inches(5.6), Inches(4.5))
        rtf_l = r_left.text_frame
        rtf_l.word_wrap = True

        rp1 = rtf_l.paragraphs[0]
        rp1.text = "Potential Risks"
        rp1.font.size = Pt(20)
        rp1.font.bold = True
        rp1.font.color.rgb = RGBColor(248, 113, 113)  # Red accent

        risks = ["LLM rate limits during live pitch", "Database connection drops", "Template placeholder misalignment"]
        for r in risks:
            rp = rtf_l.add_paragraph()
            rp.text = f"⚠ {r}"
            rp.font.size = Pt(14)
            rp.font.color.rgb = TEXT_MUTED
            rp.space_before = Pt(12)

        r_right = add_card(s9, Inches(6.9), Inches(2.2), Inches(5.6), Inches(4.5))
        rtf_r = r_right.text_frame
        rtf_r.word_wrap = True

        rsp1 = rtf_r.paragraphs[0]
        rsp1.text = "KAIROS Solution Strategy"
        rsp1.font.size = Pt(20)
        rsp1.font.bold = True
        rsp1.font.color.rgb = ACCENT_EMERALD

        solutions = ["Automated multi-model cascade (Groq -> Nvidia -> HF)", "Resilient async connection pooling", "Widescreen vector card bounding engine"]
        for sol in solutions:
            rsp = rtf_r.add_paragraph()
            rsp.text = f"✔ {sol}"
            rsp.font.size = Pt(14)
            rsp.font.color.rgb = TEXT_WHITE
            rsp.space_before = Pt(12)

        # Slide 10: Conclusion & Call to Action
        s10 = prs.slides.add_slide(blank_layout)
        add_background(s10)
        add_badge(s10, "CONCLUSION", ACCENT_PURPLE, Inches(0.8), Inches(0.6))

        c10_card = add_card(s10, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.2), ACCENT_PURPLE)
        ctf10 = c10_card.text_frame
        ctf10.word_wrap = True

        cp10_1 = ctf10.paragraphs[0]
        cp10_1.text = "Ready for Production Execution"
        cp10_1.font.size = Pt(36)
        cp10_1.font.bold = True
        cp10_1.font.color.rgb = TEXT_WHITE
        cp10_1.alignment = PP_ALIGN.CENTER

        cp10_2 = ctf10.add_paragraph()
        cp10_2.text = "KAIROS empowers hackathon teams to ship faster, present smarter, and win bigger."
        cp10_2.font.size = Pt(18)
        cp10_2.font.color.rgb = TEXT_MUTED
        cp10_2.alignment = PP_ALIGN.CENTER
        cp10_2.space_before = Pt(16)

        out_stream = io.BytesIO()
        prs.save(out_stream)
        out_stream.seek(0)
        return out_stream.getvalue()

    @staticmethod
    def fill_presentation(
        template_source: str,
        session_name: str,
        problem_statement: str,
        user_idea: str,
        pitch_sections: Dict[str, str],
        milestones: List[Dict[str, Any]] = None,
        tasks: List[Dict[str, Any]] = None,
        team_data: Dict[str, Any] = None,
        custom_pptx_bytes: bytes = None
    ) -> bytes:
        if template_source == "template-gamma" or not template_source:
            return PPTEngine.build_gamma_presentation(
                session_name=session_name,
                problem_statement=problem_statement,
                user_idea=user_idea,
                pitch_sections=pitch_sections,
                milestones=milestones,
                tasks=tasks,
                team_data=team_data
            )

        if custom_pptx_bytes:
            prs = Presentation(io.BytesIO(custom_pptx_bytes))
        else:
            t_path = PPTEngine.get_template_path(template_source)
            if not os.path.exists(t_path):
                return PPTEngine.build_gamma_presentation(
                    session_name=session_name,
                    problem_statement=problem_statement,
                    user_idea=user_idea,
                    pitch_sections=pitch_sections,
                    milestones=milestones,
                    tasks=tasks,
                    team_data=team_data
                )
            prs = Presentation(t_path)

        team_info = team_data or {}
        user_name = team_info.get('name', 'Innovator')
        user_role = team_info.get('role', 'Fullstack Engineer')
        skills = team_info.get('skills', ['Python', 'React'])
        skills_str = ', '.join(skills) if isinstance(skills, list) else str(skills)

        ms_bullets = [f"{m.get('title', 'Milestone')}: {m.get('status', 'In Progress')}" for m in (milestones or [])]
        task_bullets = [f"{t.get('name', 'Task')} [{t.get('priority', 'High')}]" for t in (tasks[:6] if tasks else [])]

        slide_contents = [
            {
                "title": session_name or "Project Pitch",
                "subtitle": user_idea or "AI Hackathon Execution Platform",
                "bullets": ["AI-Driven Co-Founder Engine", "Real-Time Task Syncing", "Zero-Overflow Slide Generation"]
            },
            {
                "title": "Problem Statement & Vision",
                "subtitle": problem_statement or "Addressing hackathon project planning and execution challenges.",
                "bullets": ["Loss of project momentum during hackathons", "Unstructured milestone management", "Manual pitch slide design overhead"]
            },
            {
                "title": "Core Solution & Product Demo",
                "subtitle": pitch_sections.get("demo", "Real-time AI workflow engine for execution teams."),
                "bullets": ["Interactive AI Coach Room", "Task Board with AI Blocker Assistance", "Instant PPTX & PDF Presentation Suite"]
            },
            {
                "title": "System Architecture & Stack",
                "subtitle": pitch_sections.get("architecture", "FastAPI backend, Supabase DB, React 19 UI."),
                "bullets": ["FastAPI Async Backend", "Supabase Database & Auth", "Claude LLM Broker & Agents"]
            },
            {
                "title": "Roadmap & Execution Plan",
                "subtitle": "Sprint Milestones",
                "bullets": ms_bullets if ms_bullets else ["Phase 1: Architecture & DB Setup", "Phase 2: AI Coach & Task Sync", "Phase 3: Presentation Studio Export"]
            },
            {
                "title": "Sprint Tasks & Progress",
                "subtitle": "Execution Metrics",
                "bullets": task_bullets if task_bullets else ["Task 1: Supabase Setup", "Task 2: AI Scope Review", "Task 3: Slide Export Engine"]
            },
            {
                "title": "Team & Technical Mastery",
                "subtitle": f"Lead: {user_name} ({user_role})",
                "bullets": [f"Lead: {user_name}", f"Role: {user_role}", f"Skills: {skills_str}"]
            },
            {
                "title": "Pitch Showcase & Impact",
                "subtitle": pitch_sections.get("showcase", "KAIROS provides an end-to-end execution co-founder."),
                "bullets": ["Reduces pitch compilation time by 90%", "Guarantees zero text overflow across all templates", "Professional PPTX and PDF output streams"]
            },
            {
                "title": "Risk Mitigation & Support",
                "subtitle": "Resilience Strategy",
                "bullets": ["LLM rate-limit fallback routing", "Resilient database connection pools", "Validated slide placeholder mapping"]
            },
            {
                "title": "Conclusion & Next Steps",
                "subtitle": "Ready for Live Execution",
                "bullets": ["Launch local servers", "Test live presentation decks", "Submit project to judges"]
            }
        ]

        footer_patterns = ["presented by", "presented to", "website :", "date", "thynk unlimited", "fradel and spies", "+123-456-7890", "reallygreatsite"]

        for slide_idx, slide in enumerate(prs.slides):
            c_data = slide_contents[slide_idx] if slide_idx < len(slide_contents) else {
                "title": f"Project Insight {slide_idx + 1}",
                "subtitle": pitch_sections.get("slides", "Comprehensive pitch execution overview."),
                "bullets": ["Key technical metric 1", "Key technical metric 2", "Key technical metric 3"]
            }

            text_shapes = [s for s in slide.shapes if s.has_text_frame]

            title_shapes = []
            body_shapes = []

            for shape in text_shapes:
                txt = shape.text_frame.text.strip().lower()

                if any(fp in txt for fp in footer_patterns):
                    if "presented by" in txt:
                        PPTEngine.fit_text_to_frame(shape.text_frame, f"Presented By : {user_name}", max_font_size=12, min_font_size=8, is_title=False)
                    elif "presented to" in txt:
                        PPTEngine.fit_text_to_frame(shape.text_frame, "Presented To : Hackathon Judges & Evaluators", max_font_size=12, min_font_size=8, is_title=False)
                    continue

                top_pos = shape.top.inches if hasattr(shape, 'top') else 0

                if (top_pos < 2.5 and len(txt) < 80) or "lorem" not in txt and len(txt) < 35 and top_pos < 3.5:
                    title_shapes.append(shape)
                else:
                    body_shapes.append(shape)

            if title_shapes:
                PPTEngine.fit_text_to_frame(title_shapes[0].text_frame, c_data["title"], max_font_size=24, min_font_size=14, is_title=True)
                if len(title_shapes) > 1 and c_data.get("subtitle"):
                    PPTEngine.fit_text_to_frame(title_shapes[1].text_frame, c_data["subtitle"], max_font_size=14, min_font_size=10, is_title=False)

            if body_shapes:
                if len(body_shapes) == 1:
                    body_text = c_data["subtitle"] + "\n\n" + "\n".join([f"• {b}" for b in c_data["bullets"]])
                    PPTEngine.fit_text_to_frame(body_shapes[0].text_frame, body_text, max_font_size=14, min_font_size=9, is_title=False)
                else:
                    bullets = c_data["bullets"]
                    for idx, b_shape in enumerate(body_shapes):
                        if idx < len(bullets):
                            card_text = bullets[idx]
                        else:
                            card_text = c_data["subtitle"]
                        PPTEngine.fit_text_to_frame(b_shape.text_frame, card_text, max_font_size=12, min_font_size=8, is_title=False)

        output_stream = io.BytesIO()
        prs.save(output_stream)
        output_stream.seek(0)
        return output_stream.getvalue()

    @staticmethod
    def generate_pdf(
        session_name: str,
        problem_statement: str,
        user_idea: str,
        pitch_sections: Dict[str, str],
        milestones: List[Dict[str, Any]] = None,
        tasks: List[Dict[str, Any]] = None
    ) -> bytes:
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(
            buffer,
            pagesize=landscape(letter),
            rightMargin=36,
            leftMargin=36,
            topMargin=36,
            bottomMargin=36
        )

        styles = getSampleStyleSheet()
        
        title_style = ParagraphStyle(
            'SlideTitle',
            parent=styles['Heading1'],
            fontSize=22,
            leading=26,
            textColor=colors.HexColor('#a855f7'),
            spaceAfter=12
        )
        
        body_style = ParagraphStyle(
            'SlideBody',
            parent=styles['BodyText'],
            fontSize=11,
            leading=16,
            textColor=colors.HexColor('#334155'),
            spaceAfter=10
        )
        
        bullet_style = ParagraphStyle(
            'SlideBullet',
            parent=body_style,
            leftIndent=15,
            firstLineIndent=-10
        )

        story = []

        # Slide 1: Cover
        story.append(Paragraph(f"<b>{session_name or 'Project Presentation'}</b>", title_style))
        story.append(Spacer(1, 15))
        story.append(Paragraph(f"<b>Vision:</b> {user_idea or 'AI-Powered Project Co-Founder & Execution Engine'}", body_style))
        story.append(Spacer(1, 20))
        story.append(Paragraph("Generated by KAIROS Pitch Studio", body_style))
        story.append(PageBreak())

        # Slide 2: Problem Statement
        story.append(Paragraph("<b>Problem Statement & Core Value</b>", title_style))
        story.append(Spacer(1, 10))
        story.append(Paragraph(problem_statement or "Addressing hackathon project planning and execution challenges.", body_style))
        story.append(PageBreak())

        # Slide 3: Demo Flow
        story.append(Paragraph("<b>Demo Flow & Core User Experience</b>", title_style))
        story.append(Spacer(1, 10))
        demo_text = pitch_sections.get("demo", "Step 1: Onboarding\nStep 2: AI Coach Roadmap\nStep 3: Execution Task Board\nStep 4: Pitch Studio Export")
        for line in demo_text.split("\n"):
            if line.strip():
                story.append(Paragraph(f"• {line.strip()}", bullet_style))
        story.append(PageBreak())

        # Slide 4: Milestones & Roadmap
        story.append(Paragraph("<b>Roadmap & Milestones</b>", title_style))
        story.append(Spacer(1, 10))
        if milestones:
            table_data = [["Milestone", "Status", "Target Output"]]
            for m in milestones:
                table_data.append([
                    m.get('title', 'Milestone'),
                    m.get('status', 'Planning'),
                    m.get('deliverable', 'Code & Demo')
                ])
            t = Table(table_data, colWidths=[250, 150, 250])
            t.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#6b21a8')),
                ('TEXTCOLOR', (0,0), (-1,0), colors.white),
                ('ALIGN', (0,0), (-1,-1), 'LEFT'),
                ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
                ('BOTTOMPADDING', (0,0), (-1,0), 8),
                ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#cbd5e1')),
            ]))
            story.append(t)
        else:
            story.append(Paragraph("1. Phase 1: MVP Architecture & Database Setup", bullet_style))
            story.append(Paragraph("2. Phase 2: AI Execution Coach Integration", bullet_style))
            story.append(Paragraph("3. Phase 3: Final Pitch Suite & Showcase", bullet_style))

        story.append(PageBreak())

        # Slide 5: Pitch Showcase Script
        story.append(Paragraph("<b>Final Pitch Script & Call to Action</b>", title_style))
        story.append(Spacer(1, 10))
        showcase_text = pitch_sections.get("showcase", "KAIROS empowers teams to turn hackathon ideas into functional products efficiently.")
        story.append(Paragraph(showcase_text, body_style))

        doc.build(story)
        buffer.seek(0)
        return buffer.getvalue()
